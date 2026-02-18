"""
Generate NAPE 2026 Competition PowerPoint - Replica of Round 1 Slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==== COLOR SCHEME (dark professional) ====
BG_DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
BG_CARD       = RGBColor(0x22, 0x22, 0x3A)   # slightly lighter card bg
ACCENT_BLUE   = RGBColor(0x00, 0x8B, 0xD0)   # bright blue accent
ACCENT_GREEN  = RGBColor(0x00, 0xC9, 0x7B)   # green accent
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)   # orange accent
ACCENT_RED    = RGBColor(0xE8, 0x4D, 0x4D)   # red accent
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)   # purple accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY      = RGBColor(0x99, 0x99, 0x99)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ==== HELPERS ====
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # small corner radius
    shape.adjustments[0] = 0.03
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=LIGHT_GRAY, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_section_header(slide, text, accent_color=ACCENT_BLUE):
    """Add a colored accent bar + title at top"""
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # title
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7), text, font_size=28, bold=True, color=WHITE)

def add_subtitle_tag(slide, text, left, top, color=ACCENT_BLUE, font_size=11):
    """Small colored tag/label"""
    tag_w = Inches(3.5)
    tag_h = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, tag_w, tag_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_card(slide, left, top, width, height, title, items, accent=ACCENT_BLUE, title_size=16, body_size=12):
    """Card with colored top border"""
    # card bg
    card = add_rect(slide, left, top, width, height, BG_CARD, border_color=RGBColor(0x33, 0x33, 0x55))
    # accent top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4), title, font_size=title_size, bold=True, color=accent)
    # bullet items
    if items:
        add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6), items, font_size=body_size, color=LIGHT_GRAY)

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Title text
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    "Strategic Evaluation of Generation\nOptions for SMU/NAPE",
    font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.5),
    "Assessing Build, Buy, and Co-Location Models for Long-Term Reliability & Growth",
    font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Team members
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
    "Team Members: Anh Bui, Minh Nguyen, Cuong Nguyen",
    font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# accent line
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.8), Inches(5.3), Inches(0.03))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# ============================================================
# SLIDE 2 — INTRO TO BUSINESS + CHALLENGES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Intro to Business")

# Left card - Business
add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(2.8),
    "Business Overview", [
        "• Power generation developer/operator",
        "• Clients: Hyperscalers (AI, cloud), data centers",
        "• 13,000 MW owned capacity (coal, gas, nuclear)",
        "• Operating in PJM RTO",
        "• Market Cap: $20B | EV: $23B | EV/EBITDA: 30x",
        "• Adj FCF/Share: $10.20 | Target: 30%+ growth",
    ], accent=ACCENT_BLUE, body_size=13)

# Right card - Challenges
add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(2.8),
    "Current Situation & Challenges", [
        "• Demand outpacing grid growth",
        "• Policy & decarbonization pressure",
        "• Capital & execution constraints",
        "• 8% cost of debt, 12% cost of equity",
        "• 60/40 debt-equity structure",
        "• BB credit rating limits flexibility",
    ], accent=ACCENT_ORANGE, body_size=13)

# Bottom - Fleet table
add_text_box(slide, Inches(0.5), Inches(4.2), Inches(5), Inches(0.4),
    "Existing Generation Fleet", font_size=16, bold=True, color=ACCENT_GREEN)

fleet_data = [
    "Nuclear Baseload: 2,200 MW  |  Gas Peaker: 3,100 MW  |  Gas Baseload: 3,400 MW",
    "Coal Intermediate: 600 MW  |  Coal Baseload: 200 MW  |  Coal/Gas Intermediate: 1,400 MW"
]
add_bullet_list(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(1.5), fleet_data, font_size=12, color=LIGHT_GRAY)

# ============================================================
# SLIDE 3 — BUILD (Financial)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Build", accent_color=ACCENT_BLUE)

add_subtitle_tag(slide, "In consideration of Market Competition & Financial Situation", Inches(0.6), Inches(0.9), color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
    "550 MW CCGT — Case Excel Assumptions", [
        "• CF ~60% → ~2.89 TWh/yr (capacity × CF × 8,760)",
        "• Gas = 40–45% of U.S. generation; backbone dispatchable fuel",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
    "Capex & Returns", [
        "• Capex = $1,150/kW → ~$633M",
        "• Merchant PJM pricing volatile; capacity auctions uncertain",
        "• Returns depend on securing long-term contracts",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.5),
    "Competitive Positioning", [
        "• Meets dispatchable needs, but weaker ESG vs nuclear/renewables",
        "• AI firms increasingly pursuing low-carbon power partnerships",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(6.8), Inches(4.3), Inches(6), Inches(2.5),
    "Industry Context", [
        "• Tariffs & supply-chain issues make renewables pricier",
        "• Higher interest rates increase cost of capital",
        "• ~12.3 GW of U.S. capacity retiring in 2025 (mostly coal/gas)",
    ], accent=ACCENT_BLUE, body_size=13)

# ============================================================
# SLIDE 4 — BUY (Financial)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Buy", accent_color=ACCENT_GREEN)

add_subtitle_tag(slide, "In consideration of Market Competition & Financial Situation", Inches(0.6), Inches(0.9), color=ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.2),
    "Immediate Capacity & Cash Flow", [
        "• Acquire existing 550 MW CCGT (~$450M modeled)",
        "• CF ~70% → ~3.37 TWh/yr, no construction delay",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.2),
    "Asset & Carbon Risk", [
        "• Older units → efficiency decay; higher O&M",
        "• Gas = ~875 lb/MWh CO₂ → policy exposure",
        "• LNG export growth → greater gas price volatility",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0),
    "Strategic Fit", [
        "• Fastest capacity addition but non-scalable",
        "• Nuclear (18–20% of U.S. supply) grows mainly through life extensions, limiting acquisition supply",
        "• Regulatory risk: future EPA rules require older thermal units to retire or install 90% carbon capture",
        "• Buy gives fast megawatts but locks into carbon-exposed, aging assets that don't scale well",
    ], accent=ACCENT_GREEN, body_size=13)

# ============================================================
# SLIDE 5 — PARTNER (Financial)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Partner", accent_color=ACCENT_ORANGE)

add_subtitle_tag(slide, "In consideration of Market Competition & Financial Situation", Inches(0.6), Inches(0.9), color=ACCENT_ORANGE)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.8),
    "Economic Strength", [
        "• Long-term offtake → stable cash flow",
        "• Shared capex improves IRR vs pure Build",
        "• Avoids full merchant PJM exposure",
    ], accent=ACCENT_ORANGE, body_size=12)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.9), Inches(2.8),
    "Market Alignment", [
        "• Data centers: 415 TWh/yr globally; U.S. = 45%",
        "• Nearly half of U.S. capacity in 5 clusters",
        "• PJM/NoVA = epicenter",
        "• AI firms already partnering to reopen nuclear",
    ], accent=ACCENT_ORANGE, body_size=12)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(2.8),
    "Competitive Advantage", [
        "• Co-located supply aligns reliability + sustainability",
        "• Most repeatable model across DC hubs",
    ], accent=ACCENT_ORANGE, body_size=12)

add_card(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5),
    "Key Takeaway", [
        "• Matches customer demand, stabilizes revenue, reduces PJM volatility exposure",
        "• Partner is the strongest option from a competitive standpoint",
    ], accent=ACCENT_ORANGE, body_size=14)

# ============================================================
# SLIDE 6 — BUILD (Execution)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Build", accent_color=ACCENT_BLUE)

add_subtitle_tag(slide, "In consideration of Operational & Execution (Resources & Supply Chain)", Inches(0.6), Inches(0.9), color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.2),
    "Execution Burden", [
        "• Industry-based build time: 24–30 months for CCGT",
        "• Permitting/supply-chain delays can extend to 22–36 months",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.2),
    "Cost Exposure", [
        "• High upfront capex; interest carry throughout construction",
        "• Full merchant risk until COD (Commercial Operating Date)",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0),
    "Operational Profile", [
        "• Mid-merit plant (~60% CF in Excel)",
        "• Higher long-term carbon cost risk vs zero-carbon tech",
        "• Build delivers a modern unit, but only for a business with strong balance sheet and high risk tolerance",
    ], accent=ACCENT_BLUE, body_size=13)

# ============================================================
# SLIDE 7 — BUY (Execution)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Buy", accent_color=ACCENT_GREEN)

add_subtitle_tag(slide, "In consideration of Operational & Execution (Resources & Supply Chain)", Inches(0.6), Inches(0.9), color=ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
    "Asset Quality Uncertainty", [
        "• Existing fleet shows mixed-age coal/gas/nuclear assets",
        "• Older assets → higher forced outages + R&M capex",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
    "Market Scarcity", [
        "• Nuclear = 18–20% of U.S. generation; growth via life extensions only",
        "• Efficient CCGTs near load hubs competitively priced",
        "• Limited availability of high-quality assets",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
    "Operational Downside", [
        "• No design/siting control",
        "• Full carbon exposure retained",
        "• Long-term operational profile is largely locked in on day one",
    ], accent=ACCENT_GREEN, body_size=13)

# ============================================================
# SLIDE 8 — PARTNER (Execution)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Partner (Co-location)", accent_color=ACCENT_ORANGE)

add_subtitle_tag(slide, "In consideration of Operational & Execution (Resources & Supply Chain)", Inches(0.6), Inches(0.9), color=ACCENT_ORANGE)

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
    "Co-located Asset Profile", [
        "• 100–200 MW blocks sized to DC load (model-driven design)",
        "• Behind-the-meter reduces interconnection delays",
        "• Case allows behind/in-front of meter",
    ], accent=ACCENT_ORANGE, body_size=13)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
    "Execution & Operations Advantages", [
        "• Long-term contracts → predictable dispatch",
        "• Smaller modular builds lower EPC exposure",
        "• AI firms already structuring JVs with power producers",
    ], accent=ACCENT_ORANGE, body_size=13)

add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
    "Complexity", [
        "• JV structure, shared capex, regulatory load",
        "• Higher upfront coordination, lower long-term risk",
        "• Revenue and offtake are contracted — operational risk is lower",
    ], accent=ACCENT_ORANGE, body_size=13)

# ============================================================
# SLIDE 9 — HOW PARTNER CAN LOOK + VERTICALS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "How Partner Can Look Like?", accent_color=ACCENT_ORANGE)

add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(2.5),
    "Goal", [
        "• 100–200 MW blocks sized to DC load",
        "• Behind-the-meter reduces interconnection delays",
        "• Case allows behind/in-front of meter",
    ], accent=ACCENT_ORANGE, body_size=13)

add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(2.5),
    "What SMU/NAPE Brings", [
        "• Expertise operating gas, solar, and nuclear assets",
        "• Deliver firm dispatchable capacity, zero-carbon baseload",
        "• 13,000 MW fleet across PJM including Northern Virginia",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(3.2),
    "Illustrative Co-location Structure", [
        "• 150 MW CCGT + 50 MW Solar serving a single AI campus",
        "• 15–20 year PPA with fixed or indexed pricing",
        "• Cost-plus or regulated-like return structure",
        "• Hyperscaler co-funds capex; SMU/NAPE owns & operates",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(3.2),
    "Vertical Suggestions", [
        "• AI & hyperscale campuses — large loads, ~12% growth/yr",
        "• Enterprise & government HPC — reliability-sensitive",
        "• Renewable-mandated hubs — states requiring low-carbon",
        "• Partner model scales well and repeats easily",
    ], accent=ACCENT_PURPLE, body_size=13)

# ============================================================
# SLIDE 10 — FORECASTING (1): Annual Energy + Capital Investment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Our Forecast and Estimation", accent_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(3.0),
    "1. Annual Energy Generation", [
        "• Scenario A (Nuclear): 6.45 TWh/year (92% CF)",
        "• Scenario B (Build Gas): 3.37 TWh/year (70% CF)",
        "• Scenario C (Acquire Gas): 2.89 TWh/year (60% CF)",
        "",
        "Nuclear provides 2.23x more annual energy than",
        "existing gas acquisition, matching world-class uptime.",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(3.0),
    "3. Total Capital Investment", [
        "• Scenario A (Nuclear): $3,000M ($3,750/kW)",
        "• Scenario B (Build Gas): $687.5M ($1,250/kW)",
        "• Scenario C (Acquire Gas): $450.0M ($818/kW)",
        "",
        "Nuclear requires 4.4x more capital but delivers",
        "zero-carbon baseload without construction delay.",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "Key Takeaway", [
        "Even with the least capital investment among all options, co-located Partner still generates the highest return",
        "on energy output — making it the most capital-efficient strategy by far.",
    ], accent=ACCENT_ORANGE, body_size=14)

# ============================================================
# SLIDE 11 — FORECASTING (2): NPV + LCOE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Our Forecast and Estimation", accent_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(3.0),
    "4. Net Present Value (NPV) at 7.68% WACC", [
        "• Scenario C (Acquire Gas): +$63.7M (Positive DCF)",
        "• Scenario B (Build Gas): -$133.6M (Value destructive)",
        "• Scenario A (Nuclear): -$1,466.8M (Base case @ $60 PPA)",
        "",
        "STRATEGIC UPSIDE: Nuclear NPV becomes positive",
        "at $68/MWh; Meta deal was $85-90; Microsoft was $100.",
    ], accent=ACCENT_GREEN, body_size=13)

add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(3.0),
    "4. LCOE Component Breakdown", [
        "• Nuclear: $65.41/MWh (Capital heavy, zero carbon)",
        "• Build Gas: $55.65/MWh (Fuel-driven, carbon exposed)",
        "• Buy Gas: $53.86/MWh (Best merchant margin)",
        "",
        "At $100/ton carbon: Gas costs rise by $36/MWh,",
        "swinging advantage permanently to Nuclear.",
    ], accent=ACCENT_BLUE, body_size=13)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "Key Insight", [
        "Co-Located Nuclear dominates — it's the only option creating substantial shareholder value.",
        "Build CCGT barely breaks even. Acquiring nuclear outright destroys value due to excessive purchase price.",
    ], accent=ACCENT_ORANGE, body_size=14)

# ============================================================
# SLIDE 12 — FORECASTING (3): FCF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Our Forecast and Estimation", accent_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0),
    "6. Free Cash Flow & Investor Benefit", [
        "• Hybrid Strategy (Partner+Buy) Adj FCF/Share: $13.80 (diluted) vs $10.20 base",
        "• FCF/Share Growth: 35.3% (vs 30% target) — Exceeds Management Goal",
        "• Market Multiple: 30x → 35-40x re-rating driven by 'AI infrastructure' story",
        "• Market Cap Potential: $20B → $44B+ (+122% value creation)",
        "• Credit Upgrade: BB to BBB potential adds $565M in debt savings (NPV)",
    ], accent=ACCENT_GREEN, body_size=14)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "Cash Flow Stability Comparison", [
        "Partner delivers the most stable and sustainable free cash flows over the entire projection period.",
        "This stability is critical for meeting the company's 30%+ Adj FCF/Share growth target.",
    ], accent=ACCENT_BLUE, body_size=14)

# ============================================================
# SLIDE 13 — NET PROFIT AFTER TAX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Our Forecast and Estimation", accent_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0),
    "8. Net Profit After Tax and Depreciation — Annual Net Income Trajectory", [
        "For steady-state annual net income (Years 7–17):",
        "• Co-Located Nuclear: ~$225M/year — Stable, contracted baseload economics",
        "• Buy Existing Gas: ~$70M/year — Limited by merchant risk and carbon pressure",
        "• Build New CCGT: ~$70M/year — Faces identical market constraints",
        "",
        "Strategy Wins: Nuclear delivers 3.2x higher annual earnings than gas alternatives.",
    ], accent=ACCENT_GREEN, body_size=14)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "Key Insight", [
        "The nuclear co-location partner option delivers 3.2x higher annual earnings than gas alternatives,",
        "driven by contracted revenues with data center partner eliminating merchant exposure.",
    ], accent=ACCENT_ORANGE, body_size=14)

# ============================================================
# SLIDE 14 — RISK ASSESSMENT (1)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Risk Assessment & Solution", accent_color=ACCENT_RED)

add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(3.0),
    "Risk 1: Dependence on single data center partner", [
        "The risk of relying on a single hyperscaler gives",
        "them too much control over pricing and revenue stability.",
        "",
        "Solutions:",
        "• Diversify partners and global locations",
        "• Use strong SLAs with clear performance & exit terms",
        "• Expand into regions like Vietnam, India, Philippines",
    ], accent=ACCENT_RED, body_size=12)

add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(3.0),
    "Risk 2: Complex revenue sharing", [
        "Unclear or complicated revenue-splits can lead",
        "to disputes and hurt margins.",
        "",
        "Solutions:",
        "• Standardize contracts into simple pricing tiers",
        "• Implement automated revenue-management tools",
        "  to apply formulas consistently",
    ], accent=ACCENT_ORANGE, body_size=12)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "Summary", [
        "Both risks are manageable through diversification, standardized contracts,",
        "and automated revenue management tools.",
    ], accent=ACCENT_BLUE, body_size=14)

# ============================================================
# SLIDE 15 — RISK ASSESSMENT (2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Risk Assessment & Solution", accent_color=ACCENT_RED)

add_card(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(3.0),
    "Risk 3: Operational misalignment", [
        "Engineering, finance, and commercial teams move",
        "at different speeds, which can slow down execution",
        "with the data-center partner.",
        "",
        "Solutions:",
        "• Align teams through shared OKRs",
        "• Run recurring cross-functional leadership meetings",
        "  to remove blockers early",
    ], accent=ACCENT_RED, body_size=12)

add_card(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(3.0),
    "Industry Facts — The Game Changers", [
        "• PJM Queue: 170,000 MW backlog; Build = 3-4 yr wait; Partner = Bypass",
        "• DC Load: Dominion Zone hit 20 GW by 2037 (3.5x forecast increase)",
        "• Nuclear CF: 90.96% fleet median; perfectly matches DC uptime",
        "• PPA Pricing: Market at $85-100/MWh vs our conservative $60/MWh",
        "• Capacity Market: Hit record $329/MW-day; adds $96M/yr to Nuclear",
    ], accent=ACCENT_PURPLE, body_size=11)

add_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
    "", [
        "Thank you for your attention!",
    ], accent=ACCENT_BLUE, body_size=20)

# ==== SAVE ====
output_path = r"c:\Users\hoang\NAPE\NAPE_Round1_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
