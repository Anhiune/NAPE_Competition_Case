"""
NAPE 2026 Case Competition — Final Presentation (15 slides)
Design: Green Business Pitch Deck palette
Fonts: Calibri (substitute for HeadingNow/Aileron/Poppins — universally available)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import os

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE — Green Business Pitch Deck
# ═══════════════════════════════════════════════════════════════
DARK_GREEN  = RGBColor(0x3A, 0x4E, 0x42)   # #3A4E42 — primary
MED_GREEN   = RGBColor(0x5E, 0x9A, 0x66)   # #5E9A66 — secondary
LIGHT_GREEN = RGBColor(0xAF, 0xD4, 0x80)   # #AFD480 — accent
BRIGHT_GREEN= RGBColor(0x6A, 0xC7, 0x58)   # #6AC758 — highlight
CREAM       = RGBColor(0xFF, 0xF4, 0xDA)   # #FFF4DA — warm cream
WARM_GRAY   = RGBColor(0xE9, 0xE3, 0xDE)   # #E9E3DE — subtle bg
CHARCOAL    = RGBColor(0x32, 0x32, 0x32)   # #323232 — text
NEAR_BLACK  = RGBColor(0x16, 0x13, 0x12)   # #161312 — headings
WARM_TAN    = RGBColor(0xB0, 0xA5, 0x99)   # #B0A599 — neutral
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)   # negative/risk
DARK_CREAM  = RGBColor(0xF1, 0xEA, 0xCD)   # #F1EACD — alt cream

# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)   # Widescreen 16:9
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=12, color=CHARCOAL, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
             font_name="Calibri"):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_bullet(text_frame, text, font_size=11, color=CHARCOAL, bold=False,
               level=0, font_name="Calibri"):
    """Add a bullet point."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = level
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    return p

def build_table(slide, rows, cols, left, top, width, height):
    """Create a table and return the table object."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_header_row(table, bg_color=DARK_GREEN, font_color=WHITE, font_size=11):
    """Style the first row as header."""
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = font_color
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

def style_cell(cell, text, font_size=10, color=CHARCOAL, bold=False,
               alignment=PP_ALIGN.LEFT, bg_color=None):
    """Set text and style for a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def add_section_label(slide, text, left, top, width=Inches(3), height=Inches(0.4)):
    """Add a small green pill-shaped section label."""
    shape = add_shape(slide, left, top, width, height, MED_GREEN)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / TEAM INTRO
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide1, DARK_GREEN)

# Decorative side bar
add_rect(slide1, Inches(0), Inches(0), Inches(0.35), SLIDE_H, MED_GREEN)

# Decorative circles (top-right accent)
for i, (x, y, r, c) in enumerate([
    (Inches(11.5), Inches(0.3), Inches(1.2), LIGHT_GREEN),
    (Inches(12.0), Inches(1.0), Inches(0.6), MED_GREEN),
]):
    shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, x, y, r, r)
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()

# Competition name
add_textbox(slide1, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
            "2026 SMU Cox NAPE Energy Innovation Case Competition",
            font_size=16, color=LIGHT_GREEN, bold=False)

# Main title
add_textbox(slide1, Inches(1.0), Inches(1.2), Inches(10.5), Inches(1.8),
            "Strategic Response to\nDissident Investor Activism",
            font_size=40, color=WHITE, bold=True)

# Subtitle
add_textbox(slide1, Inches(1.0), Inches(3.2), Inches(10), Inches(0.7),
            "JV/PPA with AWS  |  Talen Energy Precedent  |  Nuclear-Powered Data Center Strategy",
            font_size=16, color=CREAM, bold=False)

# Horizontal line
add_rect(slide1, Inches(1.0), Inches(4.1), Inches(8), Pt(3), LIGHT_GREEN)

# Team section
add_textbox(slide1, Inches(1.0), Inches(4.4), Inches(4), Inches(0.4),
            "TEAM MEMBERS", font_size=12, color=LIGHT_GREEN, bold=True)

team_members = [
    ("Anh Bui", ""),
    ("Cuong Nguyen", ""),
    ("Minh Nguyen", ""),
]
for i, (name, role) in enumerate(team_members):
    x = Inches(1.0 + i * 3.5)
    y = Inches(4.9)
    # Name pill
    pill = add_shape(slide1, x, y, Inches(3.0), Inches(0.45), MED_GREEN)
    pill.text_frame.paragraphs[0].text = name
    pill.text_frame.paragraphs[0].font.size = Pt(13)
    pill.text_frame.paragraphs[0].font.color.rgb = WHITE
    pill.text_frame.paragraphs[0].font.name = "Calibri"
    pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Date
add_textbox(slide1, Inches(1.0), Inches(6.5), Inches(5), Inches(0.4),
            "February 18, 2026  |  Finals Presentation",
            font_size=11, color=WARM_TAN, bold=False)

# Bottom accent bar
add_rect(slide1, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), BRIGHT_GREEN)

add_speaker_notes(slide1, """Welcome to our presentation for the 2026 SMU Cox NAPE Energy Innovation Case Competition.

We are presenting our strategic recommendation for responding to a dissident investor who has accumulated a 9% stake in our company — an independent power producer with 13,000 MW of capacity including 2,200 MW of nuclear generation in the PJM market.

Our recommendation: Execute a Joint Venture / Power Purchase Agreement with Amazon Web Services, modeled on Talen Energy's transformative AWS partnership.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — BACKGROUND & REAL CASE COMPARISON
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, CREAM)

# Top bar
add_rect(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide2, Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
            "Company Background & Real-World Precedent",
            font_size=26, color=WHITE, bold=True)

# Left panel: Our Company
add_section_label(slide2, "OUR COMPANY", Inches(0.5), Inches(1.2))
box1 = add_shape(slide2, Inches(0.5), Inches(1.7), Inches(5.5), Inches(3.0), WHITE, DARK_GREEN)
tf1 = box1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "Case Company Profile"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
p.font.name = "Calibri"
for item in [
    "13,000 MW total capacity (Nuclear + Gas + Coal)",
    "2,200 MW nuclear plant in PJM RTO",
    "$20B market cap  |  45M shares  |  $444/share",
    "30x EV/EBITDA  |  $10.20 Adj FCF/share",
    "9% dissident investor threatening 10% threshold",
    "Goal: 30% Adj FCF/share growth target",
    "BB credit rating  |  $3.1B net debt",
]:
    add_bullet(tf1, item, font_size=10, color=CHARCOAL)

# Right panel: Talen Energy
add_section_label(slide2, "REAL CASE: TALEN ENERGY (TLN)", Inches(6.5), Inches(1.2), width=Inches(4.5))
box2 = add_shape(slide2, Inches(6.5), Inches(1.7), Inches(6.0), Inches(3.0), DARK_GREEN)
tf2 = box2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Talen Energy — Our Blueprint"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = LIGHT_GREEN
p.font.name = "Calibri"
for item in [
    "13,100 MW total capacity — nearly identical fleet",
    "2,500 MW Susquehanna nuclear (PJM RTO)",
    "Sold Cumulus DC campus to AWS for $650M (Mar 2024)",
    "Long-term PPA: nuclear power → AWS through 2042+",
    "Stock: $60 → $389 — 6.5x appreciation",
    "Continued acquiring new DCs to expand AWS partnership",
    "Market cap: $3B → $17.6B post-deal",
]:
    add_bullet(tf2, item, font_size=10, color=WHITE)

# Bottom comparison table
tbl = build_table(slide2, 3, 5, Inches(0.5), Inches(5.0), Inches(12.0), Inches(2.0))
headers = ["Metric", "Our Company", "Talen Energy", "Similarity", "Implication"]
data_rows = [
    ["Nuclear Capacity", "2,200 MW", "2,500 MW (Susquehanna)", "88% match", "Same baseload value for DCs"],
    ["Total Fleet", "13,000 MW", "13,100 MW", "99% match", "Identical scale & fuel mix"],
]
for i, h in enumerate(headers):
    style_cell(tbl.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
for r, row in enumerate(data_rows, 1):
    for c, val in enumerate(row):
        bg = WARM_GRAY if r % 2 == 0 else WHITE
        if c == 3:
            style_cell(tbl.cell(r, c), val, font_size=10, color=DARK_GREEN, bold=True, alignment=PP_ALIGN.CENTER, bg_color=bg)
        else:
            style_cell(tbl.cell(r, c), val, font_size=10, alignment=PP_ALIGN.CENTER, bg_color=bg)

add_speaker_notes(slide2, """Our case company is a $20B independent power producer with 13,000 MW in PJM — virtually identical to Talen Energy.

Key company metrics:
- 2,200 MW nuclear, leveraged heavily in baseload
- 30x EV/EBITDA, $10.20 FCF/share, BB credit rating
- A 9% dissident investor is threatening to reach 10%, which triggers a special meeting

Talen Energy is our real-world model:
- Nearly identical fleet (13,100 MW, 2,500 MW nuclear, PJM)
- Sold Cumulus data center campus to AWS for $650M in March 2024
- Stock went from $60 to $389 — 6.5x appreciation
- Market cap grew from $3B to $17.6B
- Has continued acquiring additional data centers to expand the AWS partnership

The parallels are striking: same market, same fuel mix, same scale. If Talen can do it, so can we.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — OPTIONS 1 & 2: PROS/CONS (MARKET & FINANCIAL)
# ═══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, CREAM)

add_rect(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide3, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Market & Financial Analysis — Options 1 & 2",
            font_size=26, color=WHITE, bold=True)
add_section_label(slide3, "WHY WE PASSED ON THESE", Inches(9.5), Inches(0.25), width=Inches(3.5), height=Inches(0.4))

# Option 1 box
add_textbox(slide3, Inches(0.5), Inches(1.1), Inches(6), Inches(0.5),
            "Option 1: Acquire Data Center ($4–6B)", font_size=16, color=DARK_GREEN, bold=True)

tbl1 = build_table(slide3, 6, 2, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.6))
style_cell(tbl1.cell(0, 0), "✓  PROS", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=MED_GREEN)
style_cell(tbl1.cell(0, 1), "✗  CONS", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=RED)
pros1 = [
    "Full value chain control",
    "Multiple expansion potential",
    "Direct AI/DC exposure",
    "Revenue diversification",
    "Strategic optionality",
]
cons1 = [
    "$4–6B cost → massive debt or dilution",
    "Dilutive to FCF/share in Years 1–2",
    "No IPP has successfully done this",
    "Integration risk (power ≠ data centers)",
    "BB rating at risk → potential downgrade",
]
for i, (pro, con) in enumerate(zip(pros1, cons1)):
    bg = WARM_GRAY if i % 2 == 0 else WHITE
    style_cell(tbl1.cell(i+1, 0), pro, font_size=9, bg_color=bg)
    style_cell(tbl1.cell(i+1, 1), con, font_size=9, color=RED, bg_color=bg)

# Option 2 box
add_textbox(slide3, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5),
            "Option 2: Sell to Shell/ExxonMobil", font_size=16, color=DARK_GREEN, bold=True)

tbl2 = build_table(slide3, 6, 2, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.6))
style_cell(tbl2.cell(0, 0), "✓  PROS", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=MED_GREEN)
style_cell(tbl2.cell(0, 1), "✗  CONS", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=RED)
pros2 = [
    "Immediate 20–30% premium ($89–134/sh)",
    "Satisfies dissident with full exit",
    "Certain value realization",
    "De-risks shareholders",
    "No execution risk for shareholders",
]
cons2 = [
    "Company ceases to exist",
    "0/5 FCF growth — no future upside",
    "ESG funds divest (oil major owner)",
    "Hyperscalers may terminate PPAs",
    "NRC license transfer: 12–18 months, $10–20M",
]
for i, (pro, con) in enumerate(zip(pros2, cons2)):
    bg = WARM_GRAY if i % 2 == 0 else WHITE
    style_cell(tbl2.cell(i+1, 0), pro, font_size=9, bg_color=bg)
    style_cell(tbl2.cell(i+1, 1), con, font_size=9, color=RED, bg_color=bg)

# Bottom verdict bar
verdict_bar = add_shape(slide3, Inches(0.5), Inches(4.6), Inches(12.0), Inches(0.65), DARK_GREEN)
tf_v = verdict_bar.text_frame
tf_v.paragraphs[0].text = "FINANCIAL VERDICT: Option 1 is too expensive ($4–6B) and dilutive. Option 2 destroys long-term value. Neither unlocks the AI premium."
tf_v.paragraphs[0].font.size = Pt(12)
tf_v.paragraphs[0].font.color.rgb = WHITE
tf_v.paragraphs[0].font.bold = True
tf_v.paragraphs[0].font.name = "Calibri"
tf_v.paragraphs[0].alignment = PP_ALIGN.CENTER

# Key financial metrics row
metrics_data = [
    ("Option 1 Cost", "$4–6B", "Massive debt/dilution"),
    ("Option 2 Premium", "20–30%", "One-time, no upside"),
    ("Option 3 Cost", "$200–400M", "Funded from cash"),
    ("Option 3 Uplift", "+$513M/yr", "Immediate FCF boost"),
]
for i, (label, value, desc) in enumerate(metrics_data):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide3, x, Inches(5.5), Inches(2.85), Inches(1.6), WHITE, MED_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = MED_GREEN
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, value, font_size=22, color=DARK_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, font_size=9, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide3, """From a financial perspective, we evaluated all three options:

Option 1 — Acquire a Data Center ($4–6B):
The cost is prohibitive. At $4–6B, we'd need massive debt issuance (threatening our BB rating) or significant equity dilution. FCF/share would be dilutive for Years 1–2. No IPP has ever successfully integrated downstream into data center operations — these are fundamentally different businesses.

Option 2 — Sell to Shell/ExxonMobil:
While shareholders get an immediate 20–30% premium ($89–134/share), the company ceases to exist. Zero future upside from the AI/data center megatrend. ESG-focused investors would divest. Hyperscalers like AWS and Microsoft may terminate PPAs due to oil major ownership — destroying the very value the buyer is paying for.

Neither option unlocks the true AI premium our nuclear assets command. That's why we recommend Option 3.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — OPTION 3 (CHOSEN) — FINANCIAL CASE
# ═══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, CREAM)

add_rect(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide4, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Option 3: JV/PPA with AWS — Financial Impact",
            font_size=26, color=WHITE, bold=True)
add_section_label(slide4, "RECOMMENDED", Inches(10.5), Inches(0.25), width=Inches(2.5), height=Inches(0.4))

# Left: Key deal terms
terms_box = add_shape(slide4, Inches(0.5), Inches(1.2), Inches(5.5), Inches(4.5), WHITE, DARK_GREEN)
tf_terms = terms_box.text_frame
tf_terms.word_wrap = True
p = tf_terms.paragraphs[0]
p.text = "Deal Structure — Talen AWS Playbook"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
p.font.name = "Calibri"

for label, value in [
    ("PPA Partner:", "Amazon Web Services (AWS)"),
    ("PPA Capacity:", "1,500–2,000 MW nuclear output"),
    ("PPA Price:", "$90–100/MWh (vs $51 merchant)"),
    ("PPA Duration:", "20 years (2026–2046) + extension"),
    ("DC Campus:", "100–200 MW IT, sale to AWS for $750–900M"),
    ("Total Investment:", "$200–400M (funded from $260M cash)"),
    ("Equity Dilution:", "ZERO — no new shares issued"),
]:
    p2 = tf_terms.add_paragraph()
    run1 = p2.add_run()
    run1.text = label + " "
    run1.font.size = Pt(10)
    run1.font.bold = True
    run1.font.color.rgb = MED_GREEN
    run1.font.name = "Calibri"
    run2 = p2.add_run()
    run2.text = value
    run2.font.size = Pt(10)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = "Calibri"
    p2.space_before = Pt(4)

# Right: Financial impact table
tbl3 = build_table(slide4, 7, 3, Inches(6.3), Inches(1.2), Inches(6.2), Inches(3.2))
headers3 = ["Metric", "Pre-Deal", "Post-Deal (Yr 1)"]
for i, h in enumerate(headers3):
    style_cell(tbl3.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)

fin_data = [
    ["Nuclear Revenue", "$637M (merchant)", "$1,150M (PPA)  ▲"],
    ["Incremental Revenue", "—", "+$513M/yr  ▲"],
    ["Adj FCF/Share", "$10.20", "$16.87  (+65%)  ▲"],
    ["Market Cap", "$20B", "$30–35B  ▲"],
    ["Share Price", "$444", "$667–778  ▲"],
    ["Credit Rating", "BB", "BB+ (upgrade path)  ▲"],
]
for r, row in enumerate(fin_data):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    style_cell(tbl3.cell(r+1, 0), row[0], font_size=10, bold=True, bg_color=bg, alignment=PP_ALIGN.LEFT)
    style_cell(tbl3.cell(r+1, 1), row[1], font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
    style_cell(tbl3.cell(r+1, 2), row[2], font_size=10, bold=True, color=DARK_GREEN, bg_color=bg, alignment=PP_ALIGN.CENTER)

# Why it wins bar
why_bar = add_shape(slide4, Inches(6.3), Inches(4.6), Inches(6.2), Inches(1.0), MED_GREEN)
tf_why = why_bar.text_frame
tf_why.word_wrap = True
p = tf_why.paragraphs[0]
p.text = "Why Pros Outweigh Cons:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
for txt in [
    "✓ +65% FCF growth vs 30% target — immediately satisfies dissident",
    "✓ $200M cost vs $4–6B (Option 1) — minimal balance sheet risk",
    "✓ Preserves independence vs total sale (Option 2)",
    "✓ Proven by Talen: 6.5x stock appreciation with identical strategy",
]:
    add_bullet(tf_why, txt, font_size=9, color=WHITE, bold=False)

# Bottom: 5-year outlook
for i, (yr, mc, sp) in enumerate([
    ("Year 1", "$32B", "$711"),
    ("Year 3", "$44B", "$978"),
    ("Year 5", "$52B", "$1,156"),
]):
    x = Inches(0.5 + i * 2.0)
    card = add_shape(slide4, x, Inches(6.0), Inches(1.7), Inches(1.2), DARK_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = yr
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GREEN
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, mc, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, sp + "/share", font_size=9, color=CREAM, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide4, """Here's the financial case for Option 3 — our recommended strategy.

We propose a Joint Venture / Power Purchase Agreement with Amazon Web Services, modeled directly on Talen Energy's March 2024 deal.

Key terms:
- 20-year PPA at $90-100/MWh for 1,500-2,000 MW of nuclear output
- Current merchant price is only $51/MWh — this is an 86% revenue premium
- Data center campus: 100-200 MW IT capacity, sold/JV'd with AWS for $750-900M
- Total company investment: only $200-400M, funded from existing $260M cash + project debt
- ZERO equity dilution

Financial impact in Year 1:
- Nuclear revenue jumps from $637M to $1,150M — that's +$513M in incremental revenue
- Adj FCF/share goes from $10.20 to $16.87 — a 65% increase
- Market cap grows from $20B to $30-35B
- Share price: $444 to $667-778

By Year 5: Market cap reaches $52B, share price exceeds $1,156.

This is the only option that simultaneously delivers massive FCF growth, preserves independence, requires no dilution, and is proven by real-world precedent.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — CHART: REVENUE IMPACT & FCF TRAJECTORY
# ═══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, CREAM)

add_rect(slide5, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide5, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Option 3: Revenue & Share Price Trajectory",
            font_size=26, color=WHITE, bold=True)

# Chart 1: Revenue comparison (bar chart)
chart_data1 = CategoryChartData()
chart_data1.categories = ['Current\n(Merchant)', 'Year 1\n(PPA)', 'Year 3', 'Year 5', 'Year 10']
chart_data1.add_series('Nuclear Revenue ($M)', (637, 1150, 1200, 1270, 1550))

chart_frame1 = slide5.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2),
    Inches(5.8), Inches(3.5), chart_data1
)
chart1 = chart_frame1.chart
chart1.has_legend = False
plot1 = chart1.plots[0]
series1 = plot1.series[0]
series1.format.fill.solid()
series1.format.fill.fore_color.rgb = DARK_GREEN

# Data labels
series1.has_data_labels = True
data_labels1 = series1.data_labels
data_labels1.font.size = Pt(10)
data_labels1.font.bold = True
data_labels1.font.color.rgb = DARK_GREEN
data_labels1.number_format = '$#,##0"M"'
data_labels1.show_value = True
data_labels1.label_position = XL_LABEL_POSITION.OUTSIDE_END

chart1.category_axis.tick_labels.font.size = Pt(9)
chart1.value_axis.has_title = False
chart1.value_axis.visible = False
chart1.value_axis.major_gridlines.format.line.color.rgb = WARM_GRAY

add_textbox(slide5, Inches(0.5), Inches(4.7), Inches(5.8), Inches(0.5),
            "Nuclear Revenue: $51/MWh Merchant → $95/MWh PPA (+2%/yr escalator)",
            font_size=10, color=MED_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Chart 2: Share price trajectory
chart_data2 = CategoryChartData()
chart_data2.categories = ['Current', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5']
chart_data2.add_series('Share Price ($)', (444, 711, 844, 978, 1067, 1156))

chart_frame2 = slide5.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS, Inches(6.8), Inches(1.2),
    Inches(5.8), Inches(3.5), chart_data2
)
chart2 = chart_frame2.chart
chart2.has_legend = False
plot2 = chart2.plots[0]
series2 = plot2.series[0]
series2.format.line.color.rgb = DARK_GREEN
series2.format.line.width = Pt(3)
from pptx.enum.chart import XL_MARKER_STYLE
series2.marker.style = XL_MARKER_STYLE.CIRCLE
series2.marker.size = 10
series2.marker.format.fill.solid()
series2.marker.format.fill.fore_color.rgb = BRIGHT_GREEN

series2.has_data_labels = True
data_labels2 = series2.data_labels
data_labels2.font.size = Pt(10)
data_labels2.font.bold = True
data_labels2.font.color.rgb = DARK_GREEN
data_labels2.number_format = '$#,##0'
data_labels2.show_value = True
data_labels2.label_position = XL_LABEL_POSITION.ABOVE

chart2.category_axis.tick_labels.font.size = Pt(9)
chart2.value_axis.visible = False
chart2.value_axis.major_gridlines.format.line.color.rgb = WARM_GRAY

add_textbox(slide5, Inches(6.8), Inches(4.7), Inches(5.8), Inches(0.5),
            "Projected Share Price: $444 → $1,156 (+160% over 5 years)",
            font_size=10, color=MED_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Key insight cards at bottom
for i, (label, val) in enumerate([
    ("+$513M/yr", "Incremental Revenue"),
    ("+65%", "Year 1 FCF Growth"),
    ("$1,156/sh", "Year 5 Target Price"),
    ("0", "Shares Diluted"),
]):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide5, x, Inches(5.4), Inches(2.85), Inches(1.7), DARK_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GREEN
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, label, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide5, """These two charts show the financial trajectory of Option 3.

Left chart — Nuclear Revenue:
- Currently at $637M from merchant pricing at $51/MWh
- PPA reprices to $95/MWh, jumping to $1,150M in Year 1
- With 2% annual escalator, reaches $1,270M by Year 5 and $1,550M by Year 10
- This is $513M in incremental annual revenue from Day 1

Right chart — Share Price:
- Starting at $444/share today
- Year 1: $711 (based on $32B market cap / 45M shares)
- Year 5: $1,156 — a 160% appreciation
- Driven by EBITDA growth at 35x target multiple

Key metrics at bottom: +$513M/yr revenue, +65% FCF growth, $1,156 target, zero dilution.

This is transformative value creation with minimal risk.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — PLACEHOLDER: TEAMMATE'S TALEN FINANCIAL TIME SERIES
# ═══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, CREAM)

add_rect(slide6, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide6, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Talen Energy — Financial Transformation (Real Case)",
            font_size=26, color=WHITE, bold=True)
add_section_label(slide6, "REAL-WORLD VALIDATION", Inches(9.5), Inches(0.25), width=Inches(3.5))

# Placeholder for teammate's chart
placeholder = add_shape(slide6, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5), WHITE, MED_GREEN)
tf_ph = placeholder.text_frame
tf_ph.word_wrap = True
p = tf_ph.paragraphs[0]
p.text = "[TEAMMATE'S CHART — TALEN ENERGY FINANCIAL TIME SERIES]"
p.font.size = Pt(18)
p.font.color.rgb = MED_GREEN
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_para(tf_ph, "\nInsert time series analysis showing:", font_size=12, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph, "• Talen stock price trajectory ($60 → $389)", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph, "• Revenue/EBITDA growth post-AWS deal", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph, "• Market cap expansion ($3B → $17.6B)", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph, "• Our company's projected financial trajectory overlay", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

# Key stats sidebar
for i, (value, label) in enumerate([
    ("6.5x", "Stock Appreciation"),
    ("$650M", "Cumulus Sale to AWS"),
    ("+$998M", "Net Income Turnaround"),
    ("$17.6B", "Current Market Cap"),
    ("487%", "Market Cap Growth"),
]):
    y = Inches(1.2 + i * 1.1)
    card = add_shape(slide6, Inches(9.0), y, Inches(3.8), Inches(0.95), DARK_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.color.rgb = BRIGHT_GREEN
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, label, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide6, """This slide shows our teammate's time series analysis of Talen Energy's financial transformation.

Key data points:
- Stock price: Went from $60 (OTC post-bankruptcy May 2023) to $389 on NASDAQ — 6.5x appreciation
- The inflection point was the March 2024 AWS Cumulus deal announcement
- Revenue normalized from $3,089M to $2,115M post-restructuring, but profitability surged
- Net income went from -$1,289M to +$998M
- Market cap grew from ~$3B to $17.6B — 487% growth

This is the real-world validation of our proposed strategy. Talen executed the exact playbook we're recommending, with nearly identical assets, and created exceptional shareholder value.

Our projected trajectory overlays show similar or better potential given our stronger starting position ($20B market cap, no bankruptcy history, established credit).""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — OPTIONS 1 & 2: OPERATIONAL & EXECUTION PROS/CONS
# ═══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, CREAM)

add_rect(slide7, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide7, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Operational & Execution — Options 1 & 2",
            font_size=26, color=WHITE, bold=True)
add_section_label(slide7, "RESOURCES & SUPPLY CHAIN", Inches(9.5), Inches(0.25), width=Inches(3.5))

# Option 1 execution
add_textbox(slide7, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
            "Option 1: Acquire DC — Execution Challenges", font_size=14, color=DARK_GREEN, bold=True)

tbl_e1 = build_table(slide7, 5, 3, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.3))
for i, h in enumerate(["Challenge", "Severity", "Detail"]):
    style_cell(tbl_e1.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
exec1 = [
    ["Integration complexity", "HIGH", "Power gen ≠ DC operations — different teams, skills, culture"],
    ["Capital deployment", "HIGH", "$4–6B requires 18–24 month fundraising"],
    ["Talent acquisition", "MEDIUM", "Need DC engineers, not power plant operators"],
    ["Regulatory approvals", "MEDIUM", "FERC, PUC, antitrust reviews: 6–12 months"],
]
for r, row in enumerate(exec1):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    style_cell(tbl_e1.cell(r+1, 0), row[0], font_size=9, bold=True, bg_color=bg)
    sev_color = RED if row[1] == "HIGH" else MED_GREEN
    style_cell(tbl_e1.cell(r+1, 1), row[1], font_size=9, bold=True, color=sev_color, alignment=PP_ALIGN.CENTER, bg_color=bg)
    style_cell(tbl_e1.cell(r+1, 2), row[2], font_size=9, bg_color=bg)

# Option 2 execution
add_textbox(slide7, Inches(6.8), Inches(1.1), Inches(6), Inches(0.4),
            "Option 2: Sale — Execution Risks", font_size=14, color=DARK_GREEN, bold=True)

tbl_e2 = build_table(slide7, 5, 3, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.3))
for i, h in enumerate(["Challenge", "Severity", "Detail"]):
    style_cell(tbl_e2.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
exec2 = [
    ["NRC license transfer", "HIGH", "12–18 months, $10–20M legal costs"],
    ["ESG backlash", "HIGH", "Oil major ownership → ESG funds divest"],
    ["PPA termination risk", "HIGH", "AWS/Microsoft may exit nuclear PPAs"],
    ["Employee/ops disruption", "MEDIUM", "Cultural clash with oil major parent"],
]
for r, row in enumerate(exec2):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    style_cell(tbl_e2.cell(r+1, 0), row[0], font_size=9, bold=True, bg_color=bg)
    sev_color = RED if row[1] == "HIGH" else MED_GREEN
    style_cell(tbl_e2.cell(r+1, 1), row[1], font_size=9, bold=True, color=sev_color, alignment=PP_ALIGN.CENTER, bg_color=bg)
    style_cell(tbl_e2.cell(r+1, 2), row[2], font_size=9, bg_color=bg)

# Execution timeline comparison
timeline_bar = add_shape(slide7, Inches(0.5), Inches(4.2), Inches(12.0), Inches(0.5), DARK_GREEN)
tf_tl = timeline_bar.text_frame
tf_tl.paragraphs[0].text = "EXECUTION VERDICT: Option 1 takes 18–24 months and requires unproven skills. Option 2 faces 12–18 month NRC review. Option 3 closes in 1–3 months."
tf_tl.paragraphs[0].font.size = Pt(11)
tf_tl.paragraphs[0].font.color.rgb = WHITE
tf_tl.paragraphs[0].font.bold = True
tf_tl.paragraphs[0].font.name = "Calibri"
tf_tl.paragraphs[0].alignment = PP_ALIGN.CENTER

# Timeline visual
for i, (opt, time, color, desc) in enumerate([
    ("Option 1", "18–24 months", RED, "Fundraising → Acquisition → Integration"),
    ("Option 2", "12–18 months", RED, "NRC Review → FERC → Antitrust → Close"),
    ("Option 3", "1–3 months", BRIGHT_GREEN, "Negotiate PPA → Sign → Revenue Starts"),
]):
    y = Inches(5.0 + i * 0.8)
    label_box = add_shape(slide7, Inches(0.5), y, Inches(1.5), Inches(0.6), DARK_GREEN)
    label_box.text_frame.paragraphs[0].text = opt
    label_box.text_frame.paragraphs[0].font.size = Pt(10)
    label_box.text_frame.paragraphs[0].font.color.rgb = WHITE
    label_box.text_frame.paragraphs[0].font.bold = True
    label_box.text_frame.paragraphs[0].font.name = "Calibri"
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Timeline bar width proportional to time
    bar_width = Inches(8.0) if "18" in time else (Inches(6.5) if "12" in time else Inches(2.5))
    bar = add_shape(slide7, Inches(2.2), y, bar_width, Inches(0.6), color)
    tf_bar = bar.text_frame
    tf_bar.paragraphs[0].text = f"{time}  •  {desc}"
    tf_bar.paragraphs[0].font.size = Pt(9)
    tf_bar.paragraphs[0].font.color.rgb = WHITE
    tf_bar.paragraphs[0].font.bold = True
    tf_bar.paragraphs[0].font.name = "Calibri"
    tf_bar.paragraphs[0].alignment = PP_ALIGN.LEFT

add_speaker_notes(slide7, """From an operational and execution standpoint, Options 1 and 2 face serious challenges.

Option 1 — Acquiring a Data Center:
- Integration complexity is HIGH. Power generation and data center operations are fundamentally different businesses requiring different skills, teams, and cultures.
- Capital deployment of $4–6B requires an 18–24 month fundraising process
- Need to hire data center engineers — our expertise is in power plant operations

Option 2 — Sale to Oil Major:
- NRC license transfer is the biggest hurdle: 12–18 months, $10–20M in legal costs
- ESG backlash is severe — Shell/Exxon ownership triggers ESG fund divestment
- Hyperscalers like AWS and Microsoft may terminate PPAs if our nuclear plant is oil-major-owned

Option 3 by contrast:
- PPA negotiation closes in 1–3 months — it's a bilateral contract
- No license transfer needed — we retain nuclear ownership
- Campus development is 6–12 months but PPA revenue starts immediately
- We leverage AWS's own capital for DC build-out""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — OPTION 3 EXECUTION PLAN
# ═══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, CREAM)

add_rect(slide8, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide8, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Option 3: Execution Roadmap",
            font_size=26, color=WHITE, bold=True)

# Phase cards
phases = [
    ("Phase 1", "Months 0–3", "Negotiate & Sign\nAWS PPA", "$0", "PPA at $90–100/MWh\n1,500–2,000 MW nuclear\n20-year term", BRIGHT_GREEN),
    ("Phase 2", "Months 3–12", "Develop DC\nCampus", "$200–400M", "100–200 MW IT capacity\nCo-located with nuclear\nAWS funds majority", MED_GREEN),
    ("Phase 3", "Months 12–18", "Campus Sale\nto AWS", "+$750–900M", "Net cash: +$350–500M\nTalen precedent: $650M\nRetain expansion rights", DARK_GREEN),
    ("Phase 4", "Years 2–5", "Expansion\nRights", "+$500M–1B", "300–500 MW additional\nRepeatable platform\nCredit upgrade to BB+", DARK_GREEN),
    ("Phase 5", "Years 3–7", "DC Acquisition\nStrategy", "$250–530M", "Acquire existing DCs\nTalen Cumulus playbook\n$750M–1.6B value", MED_GREEN),
]

for i, (phase, timeline, title, cost, details, color) in enumerate(phases):
    x = Inches(0.3 + i * 2.6)
    # Phase header
    hdr = add_shape(slide8, x, Inches(1.2), Inches(2.35), Inches(0.65), color)
    tf_hdr = hdr.text_frame
    tf_hdr.word_wrap = True
    p = tf_hdr.paragraphs[0]
    p.text = f"{phase}: {timeline}"
    p.font.size = Pt(9)
    p.font.color.rgb = CREAM
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf_hdr, title, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, space_before=Pt(1))
    
    # Cost badge
    cost_badge = add_shape(slide8, x + Inches(0.3), Inches(1.95), Inches(1.75), Inches(0.45), CREAM if "$0" not in cost else BRIGHT_GREEN)
    tf_cb = cost_badge.text_frame
    p = tf_cb.paragraphs[0]
    p.text = cost
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GREEN if "$0" not in cost else WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Details card
    det = add_shape(slide8, x, Inches(2.55), Inches(2.35), Inches(1.6), WHITE, MED_GREEN)
    tf_det = det.text_frame
    tf_det.word_wrap = True
    for line in details.split("\n"):
        add_bullet(tf_det, line, font_size=9, color=CHARCOAL)

# Arrow connectors (simplified)
for i in range(4):
    x = Inches(0.3 + (i+1) * 2.6 - Inches(0.15).inches)
    arrow = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.5 + i * 2.6), Inches(1.45), Inches(0.25), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GREEN
    arrow.line.fill.background()

# Bottom insight
total_bar = add_shape(slide8, Inches(0.3), Inches(4.4), Inches(12.5), Inches(0.55), DARK_GREEN)
tf_total = total_bar.text_frame
tf_total.paragraphs[0].text = "TOTAL 5-YEAR VALUE CREATION: $30B+ market cap growth  |  $1,156/share target  |  Investment: $450–930M → Returns: $2.0–3.6B+ value"
tf_total.paragraphs[0].font.size = Pt(11)
tf_total.paragraphs[0].font.color.rgb = WHITE
tf_total.paragraphs[0].font.bold = True
tf_total.paragraphs[0].font.name = "Calibri"
tf_total.paragraphs[0].alignment = PP_ALIGN.CENTER

# Talen parallel section
add_textbox(slide8, Inches(0.3), Inches(5.2), Inches(12), Inches(0.4),
            "Validated by Talen Energy's Execution (Real Case)", font_size=14, color=DARK_GREEN, bold=True)

tbl_exec = build_table(slide8, 3, 4, Inches(0.3), Inches(5.7), Inches(12.5), Inches(1.4))
for i, h in enumerate(["Phase", "Talen (Actual)", "Our Plan", "Status"]):
    style_cell(tbl_exec.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
exec_compare = [
    ["PPA + DC Sale", "$650M Cumulus → AWS (Mar 2024)", "$750–900M campus → AWS", "Talen: ✓ Completed"],
    ["DC Acquisitions", "Acquired additional DCs post-deal", "$250–530M acquisition pipeline", "Talen: ✓ Expanding"],
]
for r, row in enumerate(exec_compare):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    for c, val in enumerate(row):
        style_cell(tbl_exec.cell(r+1, c), val, font_size=9, bg_color=bg, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide8, """This is our 5-phase execution roadmap.

Phase 1 (Months 0–3): Negotiate and sign the PPA with AWS. This is a bilateral contract — no regulatory approval needed for front-of-meter PPA. Cost: $0. Revenue starts flowing.

Phase 2 (Months 3–12): Develop a 100–200 MW data center campus co-located with our nuclear plant. Investment: $200–400M, but AWS funds the majority of the build-out.

Phase 3 (Months 12–18): Sell or JV the campus to AWS for $750–900M. Talen got $650M for a smaller campus. Net cash inflow: +$350–500M. We retain expansion rights.

Phase 4 (Years 2–5): Execute expansion phases. Add 300–500 MW of additional co-located capacity. Each phase generates incremental PPA revenue. Credit upgrade path: BB to BB+.

Phase 5 (Years 3–7): Following Talen's Cumulus Growth playbook, acquire existing data centers near our power plants. $250–530M investment for $750M–1.6B in value creation.

Total 5-year value: $30B+ market cap growth, $1,156/share target.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — CHART: EXECUTION DETAIL
# ═══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, CREAM)

add_rect(slide9, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide9, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Sensitivity Analysis — What Drives Value",
            font_size=26, color=WHITE, bold=True)

# Tornado chart
chart_data_t = CategoryChartData()
params = [
    "PPA Price ($/MWh)",
    "EV/EBITDA Multiple",
    "Nuclear Cap Factor",
    "WACC",
    "PPA Term (years)",
    "Gas Curve ($/MMBtu)",
]
chart_data_t.categories = params
chart_data_t.add_series('Low Case', (-6.2, -4.5, -3.1, -1.8, -2.0, 1.2))
chart_data_t.add_series('High Case', (4.8, 3.2, 1.8, 2.4, 1.4, -1.5))

chart_frame_t = slide9.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2),
    Inches(7.5), Inches(4.5), chart_data_t
)
chart_t = chart_frame_t.chart
plot_t = chart_t.plots[0]
plot_t.gap_width = 80

series_low = plot_t.series[0]
series_high = plot_t.series[1]
series_low.format.fill.solid()
series_low.format.fill.fore_color.rgb = RED
series_high.format.fill.solid()
series_high.format.fill.fore_color.rgb = BRIGHT_GREEN

chart_t.has_legend = True
chart_t.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_t.legend.font.size = Pt(9)

chart_t.category_axis.tick_labels.font.size = Pt(9)
chart_t.category_axis.tick_labels.font.bold = True
chart_t.value_axis.has_title = True
chart_t.value_axis.axis_title.text_frame.paragraphs[0].text = "NPV Impact ($B)"
chart_t.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)

# Key takeaways
add_textbox(slide9, Inches(8.5), Inches(1.2), Inches(4.3), Inches(0.4),
            "Key Takeaways", font_size=14, color=DARK_GREEN, bold=True)

insights = [
    ("PPA Price = #1 Driver", "$11B NPV swing ($75–$110/MWh). Run competitive bidding among AWS, Microsoft, Google."),
    ("Multiple Re-Rating", "$7.7B swing. Talen precedent (28x) is our floor; 40x is achievable."),
    ("Nuclear CF is Stable", "92% CF is industry-proven. Lowest risk factor — nuclear is reliable."),
    ("WACC Sensitivity", "Contracted PPA revenue reduces WACC over time (risk profile improves)."),
]
for i, (title, desc) in enumerate(insights):
    y = Inches(1.7 + i * 1.2)
    card = add_shape(slide9, Inches(8.5), y, Inches(4.3), Inches(1.0), WHITE, MED_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.font.name = "Calibri"
    add_para(tf, desc, font_size=8, color=CHARCOAL)

# Assumptions box
assumptions_box = add_shape(slide9, Inches(0.5), Inches(5.9), Inches(12.0), Inches(1.3), DARK_GREEN)
tf_a = assumptions_box.text_frame
tf_a.word_wrap = True
p = tf_a.paragraphs[0]
p.text = "Key Assumptions: WACC 7.68% (60/40 D/E, 8% CoD, 12% CoE)  |  PPA Base: $95/MWh  |  Term: 20 yrs  |  Nuclear CF: 92%  |  Tax: 40%  |  Escalator: 2%/yr"
p.font.size = Pt(9)
p.font.color.rgb = CREAM
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
add_para(tf_a, "NPV Range: Low Case $5.8B  →  Base Case $12.0B  →  High Case $16.8B  |  Sources: Case PDF/Excel, Talen–AWS deal terms, EIA AEO 2024", font_size=9, color=LIGHT_GREEN, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide9, """This tornado chart shows which assumptions drive Option 3's value the most.

Top drivers:
1. PPA Price is #1 — swinging $11B in NPV between $75 and $110/MWh. This is why we recommend running a competitive bidding process among AWS, Microsoft, and Google to maximize price.

2. EV/EBITDA multiple is #2 — $7.7B swing. Talen trades at 28x, which is our floor. With a proven PPA, 35-40x is achievable.

3. Nuclear capacity factor is #3 — but at 92%, it's among the most stable inputs. Industry-wide nuclear CF has been above 90% since 2010.

4. Gas curve actually helps us in both directions — higher gas prices increase wholesale power prices for our un-contracted capacity.

Bottom line: Even in the low case ($75/MWh PPA), Option 3 generates $5.8B in NPV. The base case is $12.0B. The risk-reward is overwhelmingly positive.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — PLACEHOLDER: TEAMMATE'S TALEN EXECUTION TIME SERIES
# ═══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, CREAM)

add_rect(slide10, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide10, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Talen Energy — Execution Timeline (Real Case)",
            font_size=26, color=WHITE, bold=True)
add_section_label(slide10, "REAL-WORLD EXECUTION", Inches(9.5), Inches(0.25), width=Inches(3.5))

# Placeholder
placeholder2 = add_shape(slide10, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5), WHITE, MED_GREEN)
tf_ph2 = placeholder2.text_frame
tf_ph2.word_wrap = True
p = tf_ph2.paragraphs[0]
p.text = "[TEAMMATE'S CHART — TALEN EXECUTION TIMELINE]"
p.font.size = Pt(18)
p.font.color.rgb = MED_GREEN
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_para(tf_ph2, "\nInsert execution analysis showing:", font_size=12, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph2, "• Talen's AWS deal timeline (negotiation → close → expansion)", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph2, "• Cumulus campus development milestones", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph2, "• Post-deal DC acquisitions (Cumulus Growth)", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
add_para(tf_ph2, "• Our proposed execution timeline overlay", font_size=11, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

# Key execution milestones sidebar
for i, (date, event) in enumerate([
    ("Mar 2024", "Cumulus sold to AWS ($650M)"),
    ("Jul 2024", "NASDAQ listing achieved"),
    ("2024–25", "Additional DC acquisitions"),
    ("Through 2042+", "Nuclear PPA with AWS"),
    ("Ongoing", "Cumulus Growth expansion"),
]):
    y = Inches(1.2 + i * 1.1)
    card = add_shape(slide10, Inches(9.0), y, Inches(3.8), Inches(0.95), DARK_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(14)
    p.font.color.rgb = BRIGHT_GREEN
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, event, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide10, """This slide shows our teammate's analysis of Talen Energy's actual execution timeline.

Key milestones:
- March 2024: Sold Cumulus data center campus to AWS for $650M
- July 2024: Uplisted from OTC to NASDAQ Global Select Market
- 2024–2025: Continued acquiring additional data center facilities through Cumulus Growth subsidiary
- PPA with AWS extends through 2042+

This validates that:
1. The deal can be structured and closed quickly (bilateral negotiation)
2. The value creation is immediate (stock re-rated on announcement)
3. The platform is expandable (Talen kept acquiring DCs post-deal)
4. The partnership deepens over time (AWS keeps growing its power needs)

Our plan mirrors this timeline but with a stronger starting position.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — COMPARATIVE ANALYSIS: ALL THREE OPTIONS
# ═══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, CREAM)

add_rect(slide11, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide11, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Comparative Analysis — All Three Options",
            font_size=26, color=WHITE, bold=True)

# Decision matrix chart
chart_data_dm = CategoryChartData()
chart_data_dm.categories = [
    'FCF/Share\nGrowth (25%)',
    'Multiple\nExpansion (20%)',
    'Execution\nRisk (20%)',
    'Shareholder\nValue (15%)',
    'Independence\n(10%)',
    'Satisfies\nDissident (10%)',
]
chart_data_dm.add_series('Opt 1: Acquire DC', (2, 3, 2, 2, 5, 3))
chart_data_dm.add_series('Opt 2: Sell', (0, 0, 3, 4, 0, 5))
chart_data_dm.add_series('Opt 3: JV/PPA (Rec)', (5, 5, 4, 5, 5, 5))

chart_frame_dm = slide11.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2),
    Inches(7.5), Inches(4.5), chart_data_dm
)
chart_dm = chart_frame_dm.chart
plot_dm = chart_dm.plots[0]
plot_dm.gap_width = 80
plot_dm.overlap = -20

s1 = plot_dm.series[0]
s2 = plot_dm.series[1]
s3 = plot_dm.series[2]
s1.format.fill.solid(); s1.format.fill.fore_color.rgb = LIGHT_GREEN
s2.format.fill.solid(); s2.format.fill.fore_color.rgb = WARM_TAN
s3.format.fill.solid(); s3.format.fill.fore_color.rgb = DARK_GREEN

chart_dm.has_legend = True
chart_dm.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_dm.legend.font.size = Pt(9)
chart_dm.category_axis.tick_labels.font.size = Pt(8)
chart_dm.value_axis.has_title = False
chart_dm.value_axis.maximum_scale = 5.5

# Summary cards
scores = [
    ("Option 1", "Acquire DC", "2.45/5", LIGHT_GREEN),
    ("Option 2", "Sell to Oil Major", "2.15/5", WARM_TAN),
    ("Option 3", "JV/PPA with AWS", "4.80/5", DARK_GREEN),
]
for i, (opt, desc, score, color) in enumerate(scores):
    x = Inches(8.5)
    y = Inches(1.2 + i * 1.8)
    card = add_shape(slide11, x, y, Inches(4.3), Inches(1.5), color)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = opt
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE if color != LIGHT_GREEN else CHARCOAL
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, score, font_size=24, color=WHITE if color != LIGHT_GREEN else CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, font_size=9, color=CREAM if color != LIGHT_GREEN else CHARCOAL, alignment=PP_ALIGN.CENTER)

# Winner bar
winner = add_shape(slide11, Inches(0.5), Inches(5.9), Inches(12.0), Inches(1.3), DARK_GREEN)
tf_w = winner.text_frame
tf_w.word_wrap = True
p = tf_w.paragraphs[0]
p.text = "RECOMMENDATION: Option 3 — JV/PPA with AWS"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
add_para(tf_w, "Scores 4.80/5 across all criteria  |  Only option delivering +65% FCF growth + independence + proven precedent  |  Talen's 6.5x stock appreciation validates the strategy", font_size=10, color=CREAM, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide11, """This slide shows our weighted decision matrix comparing all three options across six criteria.

Scoring (out of 5):

FCF/Share Growth (25% weight):
- Option 1: 2/5 — dilutive in Years 1–2
- Option 2: 0/5 — company ceases to exist
- Option 3: 5/5 — +65% Year 1 growth

Multiple Expansion (20%):
- Option 1: 3/5 — uncertain blended multiple
- Option 2: 0/5 — N/A
- Option 3: 5/5 — 30x to 35–40x (Talen achieved 28x from 5x)

Execution Risk (20%):
- Option 1: 2/5 — high complexity, unproven
- Option 2: 3/5 — NRC regulatory risk
- Option 3: 4/5 — proven by Talen precedent

Shareholder Value (15%):
- Option 1: 2/5 — neutral near-term
- Option 2: 4/5 — immediate premium
- Option 3: 5/5 — 2–3x upside

Independence (10%):
- Option 1: 5/5 — Yes
- Option 2: 0/5 — No
- Option 3: 5/5 — Yes

Satisfies Dissident (10%):
- Option 1: 3/5 — slow payoff
- Option 2: 5/5 — full exit
- Option 3: 5/5 — clear value creation

Weighted scores: Option 1: 2.45, Option 2: 2.15, Option 3: 4.80

Option 3 wins decisively.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — RISK ASSESSMENT
# ═══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, CREAM)

add_rect(slide12, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide12, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Risk Assessment & Mitigation",
            font_size=26, color=WHITE, bold=True)

# Risk table
tbl_risk = build_table(slide12, 7, 4, Inches(0.5), Inches(1.2), Inches(12.0), Inches(4.5))
for i, h in enumerate(["Risk", "Probability", "Impact", "Mitigation Strategy"]):
    style_cell(tbl_risk.cell(0, i), h, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)

risks = [
    ["PPA negotiation fails", "Low (20%)", "High", "Competitive bidding: AWS, Microsoft, Google, Meta all seeking nuclear power"],
    ["FERC blocks BTM arrangement", "Med (30%)", "Medium", "Dual-track: Front-of-meter PPA as primary (no FERC approval needed); BTM as backup"],
    ["Nuclear plant outage", "Low (10%)", "High", "92%+ capacity factor; dual-unit redundancy; force majeure clauses in PPA"],
    ["PPA pricing below target", "Med (25%)", "Low", "Even $75/MWh = +$290M/yr vs merchant; floor price protections in contract"],
    ["Dissident escalates before deal", "Med (35%)", "Medium", "Board seat offer; announce framework early to demonstrate progress"],
    ["DC campus construction delay", "Med (30%)", "Low", "Campus is secondary — PPA revenue starts Day 1 regardless of campus timeline"],
]
for r, row in enumerate(risks):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    style_cell(tbl_risk.cell(r+1, 0), row[0], font_size=9, bold=True, bg_color=bg)
    prob_color = DARK_GREEN if "Low" in row[1] else MED_GREEN
    style_cell(tbl_risk.cell(r+1, 1), row[1], font_size=9, bold=True, color=prob_color, alignment=PP_ALIGN.CENTER, bg_color=bg)
    imp_color = RED if row[2] == "High" else MED_GREEN
    style_cell(tbl_risk.cell(r+1, 2), row[2], font_size=9, bold=True, color=imp_color, alignment=PP_ALIGN.CENTER, bg_color=bg)
    style_cell(tbl_risk.cell(r+1, 3), row[3], font_size=9, bg_color=bg)

# Bottom insight
risk_bar = add_shape(slide12, Inches(0.5), Inches(5.9), Inches(12.0), Inches(1.3), DARK_GREEN)
tf_rb = risk_bar.text_frame
tf_rb.word_wrap = True
p = tf_rb.paragraphs[0]
p.text = "RISK PROFILE: Even in worst-case scenarios, Option 3 generates significant value"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
add_para(tf_rb, "Low case ($75/MWh PPA): +$290M/yr incremental revenue, $5.8B NPV  |  Talen faced greater risks (bankruptcy, OTC) and still succeeded  |  Our starting position is far stronger", font_size=10, color=CREAM, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide12, """Risk assessment for Option 3.

Key risks and our mitigation strategies:

1. PPA negotiation fails (20% probability, High impact):
Mitigation: Multiple hyperscalers are actively seeking nuclear power. AWS (Talen), Microsoft (Constellation/TMI), Google (Fervo), Meta (nuclear interest). A competitive process ensures we have alternatives.

2. FERC blocks behind-the-meter (30%, Medium):
Mitigation: We structure as front-of-meter PPA as primary path — no FERC approval needed. BTM is a backup for higher value capture. Talen's deal survived FERC review.

3. Nuclear outage (10%, High):
Mitigation: Industry-wide nuclear CF is 92.5%. Our plant has dual-unit redundancy. Force majeure clauses protect PPA revenue.

4. PPA pricing below target (25%, Low):
Mitigation: Even at $75/MWh (our low case), we still generate +$290M/yr incremental revenue vs merchant pricing.

5. Dissident escalation (35%, Medium):
Mitigation: Offer a board seat. Announce the deal framework early to demonstrate strategic progress.

6. Construction delay (30%, Low):
Mitigation: PPA revenue starts Day 1 regardless of campus construction timeline.

Key point: Talen faced far greater risks — emerging from bankruptcy, trading OTC, limited financial history — and still succeeded. We start from a much stronger position.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — ESG CONSIDERATION
# ═══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13, CREAM)

add_rect(slide13, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide13, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "ESG & Sustainability — Strategic Advantage",
            font_size=26, color=WHITE, bold=True)

# ESG comparison chart
chart_data_esg = CategoryChartData()
chart_data_esg.categories = ['Carbon\nIntensity', 'Regulatory\nESG Risk', 'Hyperscaler\nAlignment', 'SEC Climate\nReadiness', 'Community\nImpact', 'Long-Term\nSustainability']
chart_data_esg.add_series('Opt 1: Acquire DC', (3, 2, 4, 3, 3, 3))
chart_data_esg.add_series('Opt 2: Sell to Oil', (1, 1, 1, 2, 1, 1))
chart_data_esg.add_series('Opt 3: JV/PPA', (5, 4, 5, 5, 5, 5))

chart_frame_esg = slide13.shapes.add_chart(
    XL_CHART_TYPE.RADAR_FILLED, Inches(0.5), Inches(1.0),
    Inches(6.5), Inches(4.5), chart_data_esg
)
chart_esg = chart_frame_esg.chart
chart_esg.has_legend = True
chart_esg.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_esg.legend.font.size = Pt(8)
plot_esg = chart_esg.plots[0]
s_esg1 = plot_esg.series[0]; s_esg1.format.fill.solid(); s_esg1.format.fill.fore_color.rgb = LIGHT_GREEN
s_esg2 = plot_esg.series[1]; s_esg2.format.fill.solid(); s_esg2.format.fill.fore_color.rgb = WARM_TAN
s_esg3 = plot_esg.series[2]; s_esg3.format.fill.solid(); s_esg3.format.fill.fore_color.rgb = MED_GREEN

# Hyperscaler targets
add_textbox(slide13, Inches(7.3), Inches(1.0), Inches(5.5), Inches(0.4),
            "Hyperscaler Sustainability Mandates", font_size=13, color=DARK_GREEN, bold=True)

hyper_data = [
    ("AWS", "Net-zero by 2040  |  100% RE by 2025", "35+ GW needed"),
    ("Microsoft", "Carbon-negative by 2030  |  24/7 CFE", "30+ GW needed"),
    ("Google", "24/7 carbon-free by 2030", "25+ GW needed"),
    ("Meta", "Net-zero by 2030", "15+ GW needed"),
]
for i, (name, target, power) in enumerate(hyper_data):
    y = Inches(1.5 + i * 0.85)
    card = add_shape(slide13, Inches(7.3), y, Inches(5.5), Inches(0.75), WHITE, MED_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{name}: "
    run1.font.size = Pt(10)
    run1.font.bold = True
    run1.font.color.rgb = DARK_GREEN
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = target
    run2.font.size = Pt(9)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = "Calibri"
    add_para(tf, f"Power demand: {power}", font_size=8, color=MED_GREEN, bold=True)

# ESG value quantification
esg_bar = add_shape(slide13, Inches(0.5), Inches(5.7), Inches(12.0), Inches(1.5), DARK_GREEN)
tf_esg = esg_bar.text_frame
tf_esg.word_wrap = True
p = tf_esg.paragraphs[0]
p.text = "ESG Value Creation: $2–4B Over 5 Years"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
add_para(tf_esg, "+1–2x EV/EBITDA from ESG premium ($1.3–2.6B)  |  Green bond savings: $15–25M/yr  |  PPA premium: +$3–5/MWh ($36–60M/yr)  |  45Q CCUS credits: $85–350M/yr", font_size=9, color=CREAM, alignment=PP_ALIGN.CENTER)
add_para(tf_esg, "Option 2 DESTROYS ESG value: Oil major ownership → ESG fund divestment (15–25% of investor base lost)  |  Option 3 CREATES ESG premium", font_size=9, color=BRIGHT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide13, """ESG is not just a checkbox — it's a strategic weapon for Option 3.

The radar chart shows Option 3 dominates all ESG dimensions:
- Carbon intensity: Nuclear = 0 g CO2/kWh vs Gas 410g, Coal 820g
- Hyperscaler alignment: All four major hyperscalers have aggressive carbon-free targets
- SEC climate disclosure: Nuclear baseload + CCUS roadmap = best-in-class compliance

Hyperscaler sustainability mandates create massive demand for our offering:
- AWS: Net-zero by 2040, 35+ GW of power needed
- Microsoft: Carbon-negative by 2030, willing to pay premium for nuclear (TMI restart)
- Google: 24/7 carbon-free by 2030
- Meta: Net-zero by 2030

ESG value quantification:
- +1–2x EV/EBITDA premium from ESG investors = $1.3–2.6B market cap
- Green bond eligibility saves 30–50 bps = $15–25M/yr
- PPA premium for carbon-free power: +$3–5/MWh = $36–60M/yr
- 45Q CCUS tax credits: $85–350M/yr (phased)
- Total: $2–4B in ESG-driven value over 5 years

Critical contrast: Option 2 (sale to Shell/Exxon) would destroy ESG value. ESG funds would divest, and hyperscalers may terminate PPAs.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — CCUS INTEGRATION
# ═══════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide14, CREAM)

add_rect(slide14, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide14, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "CCUS Integration — Gas Fleet Value Unlock",
            font_size=26, color=WHITE, bold=True)

# CCUS value stack chart
chart_data_ccus = CategoryChartData()
chart_data_ccus.categories = ['45Q Tax\nCredit', 'Carbon Cost\nAvoidance', 'Premium PPA\nPricing', 'Total CCUS\nValue']
chart_data_ccus.add_series('Value ($/MWh)', (12.5, 8.0, 7.5, 28.0))

chart_frame_ccus = slide14.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2),
    Inches(5.5), Inches(3.0), chart_data_ccus
)
chart_ccus = chart_frame_ccus.chart
chart_ccus.has_legend = False
plot_ccus = chart_ccus.plots[0]
series_ccus = plot_ccus.series[0]
series_ccus.format.fill.solid()
series_ccus.format.fill.fore_color.rgb = MED_GREEN

series_ccus.has_data_labels = True
dl_ccus = series_ccus.data_labels
dl_ccus.font.size = Pt(11)
dl_ccus.font.bold = True
dl_ccus.font.color.rgb = DARK_GREEN
dl_ccus.number_format = '$#,##0.0"/MWh"'
dl_ccus.show_value = True
dl_ccus.label_position = XL_LABEL_POSITION.OUTSIDE_END

chart_ccus.category_axis.tick_labels.font.size = Pt(9)
chart_ccus.value_axis.visible = False

# Phased deployment table
tbl_ccus = build_table(slide14, 5, 5, Inches(6.3), Inches(1.2), Inches(6.5), Inches(3.0))
for i, h in enumerate(["Phase", "Capacity", "45Q Credits/yr", "Capex", "Timeline"]):
    style_cell(tbl_ccus.cell(0, i), h, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
ccus_data = [
    ["1: Pilot", "800 MW", "$85M", "$400M", "2027–29"],
    ["2: Scale", "2,000 MW", "$210M", "$900M", "2029–31"],
    ["3: Full Fleet", "3,500 MW", "$350M", "$1.5B", "2031–34"],
    ["Total", "6,300 MW", "$350M/yr", "$2.8B", "2027–34"],
]
for r, row in enumerate(ccus_data):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    if r == 3:
        bg = DARK_GREEN
    for c, val in enumerate(row):
        if r == 3:
            style_cell(tbl_ccus.cell(r+1, c), val, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=bg)
        else:
            style_cell(tbl_ccus.cell(r+1, c), val, font_size=9, alignment=PP_ALIGN.CENTER, bg_color=bg)

# Competitive moat explanation
moat = add_shape(slide14, Inches(0.5), Inches(4.5), Inches(12.0), Inches(1.5), WHITE, MED_GREEN)
tf_moat = moat.text_frame
tf_moat.word_wrap = True
p = tf_moat.paragraphs[0]
p.text = "Why CCUS Creates an Unbeatable Competitive Moat"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
p.font.name = "Calibri"

tbl_moat = build_table(slide14, 4, 3, Inches(0.7), Inches(5.2), Inches(11.6), Inches(1.4))
for i, h in enumerate(["Our Offering", "Carbon Intensity", "Availability"]):
    style_cell(tbl_moat.cell(0, i), h, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=MED_GREEN)
moat_data = [
    ["Nuclear (2,200 MW)", "0 lb CO₂/MWh", "92% CF"],
    ["Gas + CCUS (6,500 MW)", "~72 lb CO₂/MWh (90% capture)", "Dispatchable"],
    ["Combined Fleet", "~15 lb CO₂/MWh blended", "99.9% reliability"],
]
for r, row in enumerate(moat_data):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    if r == 2:
        for c, val in enumerate(row):
            style_cell(tbl_moat.cell(r+1, c), val, font_size=9, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.CENTER, bg_color=bg)
    else:
        for c, val in enumerate(row):
            style_cell(tbl_moat.cell(r+1, c), val, font_size=9, alignment=PP_ALIGN.CENTER, bg_color=bg)

# Bottom bar
ccus_bar = add_shape(slide14, Inches(0.5), Inches(6.8), Inches(12.0), Inches(0.5), DARK_GREEN)
tf_cb2 = ccus_bar.text_frame
tf_cb2.paragraphs[0].text = "24/7 NEAR-ZERO CARBON: Nuclear + CCUS Gas = Only IPP offering true 24/7 carbon-free energy in PJM"
tf_cb2.paragraphs[0].font.size = Pt(11)
tf_cb2.paragraphs[0].font.color.rgb = WHITE
tf_cb2.paragraphs[0].font.bold = True
tf_cb2.paragraphs[0].font.name = "Calibri"
tf_cb2.paragraphs[0].alignment = PP_ALIGN.CENTER

add_speaker_notes(slide14, """CCUS Integration strengthens Option 3 by addressing our 6,500 MW gas fleet.

The value stack chart shows CCUS adds $28/MWh to gas generation:
- $12.5/MWh from 45Q tax credits ($85/ton for geological storage under the IRA)
- $8.0/MWh from carbon cost avoidance (at $25/ton carbon price)
- $7.5/MWh from premium PPA pricing (carbon-free gas commands higher prices)

Phased deployment:
- Phase 1 (2027–29): 800 MW pilot on most efficient CCGT, $85M/yr in 45Q credits
- Phase 2 (2029–31): Scale to 2,000 MW baseload gas, $210M/yr
- Phase 3 (2031–34): Full fleet at 3,500 MW, $350M/yr
- Total investment: $2.8B over 7 years

The competitive moat: No other IPP in PJM can offer 24/7 near-zero carbon energy.
- Nuclear: 0 g CO2/MWh, 92% capacity factor
- Gas + CCUS: ~72 lb/MWh (90% capture), always available
- Combined: ~15 lb/MWh blended, 99.9% reliability

This makes us the only choice for hyperscalers who need true 24/7 carbon-free energy.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — NUCLEAR REGULATORY & RISK SOLUTION
# ═══════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide15, CREAM)

add_rect(slide15, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARK_GREEN)
add_textbox(slide15, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
            "Nuclear Regulatory & Risk Management",
            font_size=26, color=WHITE, bold=True)

# Regulatory matrix
tbl_reg = build_table(slide15, 5, 4, Inches(0.5), Inches(1.2), Inches(12.0), Inches(2.8))
for i, h in enumerate(["Regulatory Body", "Key Concern", "Risk Level", "Our Solution"]):
    style_cell(tbl_reg.cell(0, i), h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=DARK_GREEN)
reg_data = [
    ["NRC", "License amendment for co-located DC", "Low–Medium", "No change needed for front-of-meter PPA; BTM has Talen precedent (approved)"],
    ["FERC", "BTM load reduces grid capacity", "Medium", "Front-of-meter PPA as primary (no approval needed); Talen–AWS survived FERC"],
    ["State PUC", "Siting permits for DC campus", "Low", "Nuclear sites have industrial zoning; PA PUC approved Talen campus"],
    ["PJM RTO", "Capacity obligation if BTM", "Medium", "Maintain FoM PPA to preserve capacity revenue; replace with gas fleet"],
]
for r, row in enumerate(reg_data):
    bg = WARM_GRAY if r % 2 == 0 else WHITE
    style_cell(tbl_reg.cell(r+1, 0), row[0], font_size=10, bold=True, bg_color=bg, alignment=PP_ALIGN.CENTER)
    style_cell(tbl_reg.cell(r+1, 1), row[1], font_size=9, bg_color=bg)
    risk_color = DARK_GREEN if "Low" in row[2] else MED_GREEN
    style_cell(tbl_reg.cell(r+1, 2), row[2], font_size=10, bold=True, color=risk_color, alignment=PP_ALIGN.CENTER, bg_color=bg)
    style_cell(tbl_reg.cell(r+1, 3), row[3], font_size=9, bg_color=bg)

# Dual-track strategy visual
add_textbox(slide15, Inches(0.5), Inches(4.2), Inches(6), Inches(0.4),
            "Dual-Track Regulatory Strategy", font_size=13, color=DARK_GREEN, bold=True)

# Track A
track_a = add_shape(slide15, Inches(0.5), Inches(4.7), Inches(5.8), Inches(1.0), BRIGHT_GREEN)
tf_ta = track_a.text_frame
tf_ta.word_wrap = True
p = tf_ta.paragraphs[0]
p.text = "Track A: Front-of-Meter PPA (PRIMARY — LOW RISK)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
add_para(tf_ta, "Sells via PJM grid  →  No FERC approval  →  Preserves capacity revenue  →  Standard bilateral PPA", font_size=9, color=WHITE)

# Track B
track_b = add_shape(slide15, Inches(0.5), Inches(5.85), Inches(5.8), Inches(1.0), MED_GREEN)
tf_tb = track_b.text_frame
tf_tb.word_wrap = True
p = tf_tb.paragraphs[0]
p.text = "Track B: Behind-the-Meter (BACKUP — HIGHER VALUE)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
add_para(tf_tb, "Direct power delivery  →  NRC 10 CFR 50.90 amendment  →  Saves $10–15/MWh transmission  →  Talen precedent exists", font_size=9, color=WHITE)

# Key advantages
add_textbox(slide15, Inches(6.8), Inches(4.2), Inches(6), Inches(0.4),
            "Key Regulatory Advantages (vs Options 1 & 2)", font_size=13, color=DARK_GREEN, bold=True)

advantages = [
    ("No NRC License Transfer", "Option 2 requires full NRC transfer (12–18 months, $10–20M). Option 3: we retain ownership — no transfer needed."),
    ("No Antitrust Review", "Option 1's $4–6B acquisition triggers HSR Act review. Option 3: bilateral PPA, no concentration concerns."),
    ("FERC De-Risked", "Front-of-meter PPA is standard bilateral trade. BTM is backup only — Talen precedent exists. Minimal regulatory friction."),
    ("Speed to Close", "Option 3 closes in 1–3 months. Revenue starts flowing immediately while campus develops in parallel."),
]
for i, (title, desc) in enumerate(advantages):
    y = Inches(4.7 + i * 0.72)
    card = add_shape(slide15, Inches(6.8), y, Inches(5.8), Inches(0.65), WHITE, MED_GREEN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"✓ {title}: "
    run1.font.size = Pt(9)
    run1.font.bold = True
    run1.font.color.rgb = DARK_GREEN
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(8)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = "Calibri"

add_speaker_notes(slide15, """Nuclear regulatory considerations and our risk management approach.

Four regulatory bodies and our solutions:

1. NRC (Nuclear Regulatory Commission):
- No license amendment needed for front-of-meter PPA
- Behind-the-meter would require 10 CFR 50.90 amendment, but Talen got this approved
- Risk: Low-Medium

2. FERC:
- Behind-the-meter reduces grid capacity — this is the biggest regulatory concern
- Our solution: Front-of-meter PPA as primary track (no FERC approval needed)
- Talen-AWS survived FERC review, setting favorable precedent
- Risk: Medium

3. State PUC:
- Siting permits for DC campus is straightforward — nuclear sites already have industrial zoning
- Risk: Low

4. PJM RTO:
- Capacity obligation preserved with front-of-meter PPA
- Can offer replacement capacity from gas fleet for BTM scenario
- Risk: Medium

Dual-Track Strategy:
- Track A (Primary): Front-of-meter PPA — standard bilateral trade, no regulatory approvals needed, preserves capacity market revenue
- Track B (Backup): Behind-the-meter — higher value ($10-15/MWh transmission savings) but requires NRC amendment and faces FERC scrutiny

Key advantage: Option 3 avoids the 12-18 month NRC license transfer that Option 2 requires, and avoids the antitrust reviews that Option 1 triggers.""")

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "NAPE_Final_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
